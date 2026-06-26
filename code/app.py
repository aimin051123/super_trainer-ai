"""
SuperTutor v5.0 — 基于 DeepSeek API 的智能学习辅助系统
技术设计依据：技术设计文档5.0.md
架构：四层模型（UI → 业务逻辑 → 数据/AI → 外部服务）
数据库：6 张表（course / knowledge_point / quiz_record / quiz_question / wrong_answer / chat_history）
"""
import os, re, json, sqlite3, configparser
from datetime import datetime, timedelta
from pathlib import Path
import streamlit as st
from openai import OpenAI

# ==================== 配置 ====================
BASE_DIR = Path(__file__).parent
CFG_FILE = BASE_DIR / "config.properties"
_config = configparser.ConfigParser()
if CFG_FILE.exists():
    _config.read(str(CFG_FILE), encoding="utf-8")
BASE_URL = _config.get("DEFAULT", "DEEPSEEK_BASE_URL", fallback="https://api.deepseek.com")
MODEL = _config.get("DEFAULT", "DEEPSEEK_MODEL", fallback="deepseek-chat")
TIMEOUT = _config.getint("DEFAULT", "DEEPSEEK_TIMEOUT", fallback=30)

def get_api_key():
    """侧边栏优先，其次环境变量，最后 config.properties"""
    key = st.session_state.get("api_key", "")
    if key:
        return key
    key = os.environ.get("DEEPSEEK_API_KEY", "")
    if key:
        return key
    key = _config.get("DEFAULT", "DEEPSEEK_API_KEY", fallback="")
    if key and key != "sk-your-key-here":
        return key
    return ""

ENTITIES = ["概念", "算法", "对比", "考试技巧"]
DIFFS = ["自动匹配", "基础", "中等", "困难"]
COUNTS = [3, 5, 8, 10]
KB_ROOT = BASE_DIR / "知识库"
RAW = KB_ROOT / "原始资料"
CPD = KB_ROOT / "编译后知识"
IDX = KB_ROOT / "全局索引.md"
LOG = KB_ROOT / "变更日志.md"
DB = BASE_DIR / "mastery.db"

# ==================== 数据库 ====================
def db_query(sql, params=(), fetch=False):
    conn = sqlite3.connect(str(DB))
    conn.execute("PRAGMA foreign_keys = ON")
    cur = conn.cursor()
    cur.execute(sql, params)
    if fetch:
        rows = cur.fetchall()
        conn.close()
        return rows
    conn.commit()
    conn.close()

def db_execute_many(sql, params_list):
    """批量插入"""
    conn = sqlite3.connect(str(DB))
    conn.execute("PRAGMA foreign_keys = ON")
    conn.executemany(sql, params_list)
    conn.commit()
    conn.close()

def init_db():
    """v5.0：6 张表，含外键和级联删除"""
    db_query("""CREATE TABLE IF NOT EXISTS course (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT NOT NULL,
        description TEXT DEFAULT '',
        exam_date TEXT DEFAULT '',
        daily_study_minutes INTEGER DEFAULT 240,
        created_at TEXT DEFAULT (datetime('now','localtime'))
    )""")
    db_query("""CREATE TABLE IF NOT EXISTS knowledge_point (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        course_id INTEGER NOT NULL,
        name TEXT NOT NULL,
        description TEXT DEFAULT '',
        importance TEXT DEFAULT 'medium',
        mastery REAL DEFAULT 0.0,
        total_count INTEGER DEFAULT 0,
        correct_count INTEGER DEFAULT 0,
        is_weak INTEGER DEFAULT 0,
        weak_reason TEXT DEFAULT '',
        last_reviewed TEXT DEFAULT '',
        FOREIGN KEY (course_id) REFERENCES course(id) ON DELETE CASCADE
    )""")
    db_query("""CREATE TABLE IF NOT EXISTS quiz_record (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        course_id INTEGER NOT NULL,
        quiz_type TEXT NOT NULL DEFAULT 'periodic',
        total_questions INTEGER DEFAULT 0,
        correct_count INTEGER DEFAULT 0,
        accuracy REAL DEFAULT 0.0,
        created_at TEXT DEFAULT (datetime('now','localtime')),
        FOREIGN KEY (course_id) REFERENCES course(id) ON DELETE CASCADE
    )""")
    db_query("""CREATE TABLE IF NOT EXISTS quiz_question (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        quiz_record_id INTEGER NOT NULL,
        course_id INTEGER NOT NULL,
        knowledge_point_id INTEGER,
        question_index INTEGER DEFAULT 1,
        question_type TEXT DEFAULT 'choice',
        content TEXT NOT NULL,
        options TEXT DEFAULT '[]',
        correct_answer TEXT NOT NULL,
        user_answer TEXT DEFAULT '',
        explanation TEXT DEFAULT '',
        difficulty TEXT DEFAULT '基础',
        created_at TEXT DEFAULT (datetime('now','localtime')),
        FOREIGN KEY (quiz_record_id) REFERENCES quiz_record(id) ON DELETE CASCADE,
        FOREIGN KEY (course_id) REFERENCES course(id) ON DELETE CASCADE,
        FOREIGN KEY (knowledge_point_id) REFERENCES knowledge_point(id) ON DELETE SET NULL
    )""")
    db_query("""CREATE TABLE IF NOT EXISTS wrong_answer (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        course_id INTEGER NOT NULL,
        knowledge_point_id INTEGER,
        quiz_question_id INTEGER,
        question_type TEXT DEFAULT 'choice',
        content TEXT NOT NULL,
        options TEXT DEFAULT '[]',
        correct_answer TEXT NOT NULL,
        user_answer TEXT DEFAULT '',
        explanation TEXT DEFAULT '',
        review_count INTEGER DEFAULT 0,
        last_reviewed TEXT DEFAULT '',
        mastered INTEGER DEFAULT 0,
        created_at TEXT DEFAULT (datetime('now','localtime')),
        FOREIGN KEY (course_id) REFERENCES course(id) ON DELETE CASCADE,
        FOREIGN KEY (knowledge_point_id) REFERENCES knowledge_point(id) ON DELETE SET NULL
    )""")
    db_query("""CREATE TABLE IF NOT EXISTS chat_history (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        course_id INTEGER NOT NULL,
        role TEXT NOT NULL,
        content TEXT NOT NULL,
        timestamp TEXT DEFAULT (datetime('now','localtime')),
        FOREIGN KEY (course_id) REFERENCES course(id) ON DELETE CASCADE
    )""")

# ==================== 课程管理 ====================
def course_create(name, description=""):
    db_query("INSERT INTO course(name,description) VALUES(?,?)", (name, description))
    return db_query("SELECT last_insert_rowid()", fetch=True)[0][0]

def course_list():
    return db_query("SELECT * FROM course ORDER BY created_at DESC", fetch=True)

def course_delete(course_id):
    db_query("DELETE FROM course WHERE id=?", (course_id,))

def course_update(course_id, exam_date=None, daily_minutes=None):
    if exam_date:
        db_query("UPDATE course SET exam_date=? WHERE id=?", (str(exam_date), course_id))
    if daily_minutes:
        db_query("UPDATE course SET daily_study_minutes=? WHERE id=?", (daily_minutes, course_id))

def course_get(course_id):
    rows = db_query("SELECT * FROM course WHERE id=?", (course_id,), fetch=True)
    return rows[0] if rows else None

# ==================== 知识点操作 ====================
def kp_upsert(course_id, name, description="", importance="medium"):
    """插入或更新知识点"""
    existing = db_query("SELECT id FROM knowledge_point WHERE course_id=? AND name=?",
                        (course_id, name), fetch=True)
    if existing:
        return existing[0][0]
    db_query("INSERT INTO knowledge_point(course_id,name,description,importance) VALUES(?,?,?,?)",
             (course_id, name, description, importance))
    return db_query("SELECT last_insert_rowid()", fetch=True)[0][0]

def kp_list(course_id):
    return db_query("SELECT * FROM knowledge_point WHERE course_id=? ORDER BY name",
                    (course_id,), fetch=True)

def kp_update_mastery(kp_id, is_correct):
    """更新单个知识点的掌握度，并调用薄弱点判定"""
    row = db_query("SELECT mastery,total_count,correct_count FROM knowledge_point WHERE id=?",
                   (kp_id,), fetch=True)
    if not row:
        return
    mastery, total, correct = row[0]
    total += 1
    correct += int(is_correct)
    new_mastery = correct / total if total > 0 else 0
    is_weak, reason = judge_weak(kp_id, new_mastery, total, is_correct)
    db_query("""UPDATE knowledge_point
        SET mastery=?, total_count=?, correct_count=?, is_weak=?, weak_reason=?, last_reviewed=?
        WHERE id=?""",
        (new_mastery, total, correct, int(is_weak), reason, datetime.now().isoformat(), kp_id))

def kp_get_weak(course_id):
    return db_query("SELECT * FROM knowledge_point WHERE course_id=? AND is_weak=1",
                    (course_id,), fetch=True)

def kp_get_observing(course_id):
    return db_query("SELECT * FROM knowledge_point WHERE course_id=? AND is_weak=0 AND weak_reason='待观察'",
                    (course_id,), fetch=True)

# ==================== 错题管理 ====================
def wrong_record(course_id, kp_id, qq_id, q_type, content, options, correct_ans, user_ans, explanation):
    db_query("""INSERT INTO wrong_answer(course_id,knowledge_point_id,quiz_question_id,
        question_type,content,options,correct_answer,user_answer,explanation)
        VALUES(?,?,?,?,?,?,?,?,?)""",
        (course_id, kp_id, qq_id, q_type, content, options, correct_ans, user_ans, explanation))

def wrong_list_unmastered(course_id):
    return db_query("""SELECT * FROM wrong_answer WHERE course_id=? AND mastered=0
        ORDER BY created_at DESC""", (course_id,), fetch=True)

def wrong_list_by_kp(course_id, kp_id):
    return db_query("""SELECT * FROM wrong_answer WHERE course_id=? AND knowledge_point_id=? AND mastered=0
        ORDER BY review_count DESC""", (course_id, kp_id), fetch=True)

def wrong_review(wrong_id, user_answer, is_correct):
    row = db_query("SELECT review_count FROM wrong_answer WHERE id=?", (wrong_id,), fetch=True)
    if row:
        new_count = row[0][0] + 1
        mastered = 1 if is_correct else 0
        db_query("""UPDATE wrong_answer SET review_count=?, last_reviewed=?, mastered=?
            WHERE id=?""", (new_count, datetime.now().isoformat(), mastered, wrong_id))
        return mastered
    return 0

# ==================== 对话历史 ====================
def chat_save(course_id, role, content):
    db_query("INSERT INTO chat_history(course_id,role,content) VALUES(?,?,?)",
             (course_id, role, content))

def chat_recent(course_id, n=10):
    rows = db_query("SELECT role,content FROM chat_history WHERE course_id=? ORDER BY id DESC LIMIT ?",
                    (course_id, n), fetch=True)
    rows.reverse()
    return rows

# ==================== AI 调用 ====================
def call_ai(prompt, system_prompt=None, json_mode=False, stream=False, temperature=0.7):
    """v5.0 统一入口：支持 json_mode / stream / temperature"""
    key = get_api_key()
    if not key:
        return "⚠️ 请先设置 DEEPSEEK_API_KEY（创建 config.properties 或设置环境变量）"
    client = OpenAI(api_key=key, base_url=BASE_URL)
    msgs = []
    if system_prompt:
        msgs.append({"role": "system", "content": system_prompt})
    msgs.append({"role": "user", "content": prompt})
    kwargs = {"model": MODEL, "messages": msgs, "temperature": temperature, "max_tokens": 4096, "timeout": TIMEOUT}
    if json_mode:
        kwargs["response_format"] = {"type": "json_object"}
    if stream:
        kwargs["stream"] = True
    try:
        response = client.chat.completions.create(**kwargs)
        if stream:
            return (chunk.choices[0].delta.content or "" for chunk in response)
        return response.choices[0].message.content
    except Exception as e:
        return f"⚠️ AI调用失败：{e}"

# ==================== SSE 流式显示 ====================
def stream_display(prompt, system_prompt=None):
    """批量刷新：每 3 token 或遇到标点才更新 UI"""
    placeholder = st.empty()
    full_text = ""
    buffer = []
    for token in call_ai(prompt, system_prompt, stream=True):
        if isinstance(token, str) and token.startswith("⚠️"):
            placeholder.error(token)
            return token
        buffer.append(token)
        full_text += token
        if len(buffer) >= 3 or token in "，。！？\n":
            placeholder.markdown(full_text + "▌")
            buffer = []
    placeholder.markdown(full_text)
    return full_text

# ==================== 文件解析 ====================
def parse_file(uploaded_file):
    name = uploaded_file.name.lower()
    ext = Path(name).suffix
    if ext == ".pdf":
        from pypdf import PdfReader
        return "\n".join(p.extract_text() or "" for p in PdfReader(uploaded_file).pages)
    if ext in (".ppt", ".pptx"):
        from pptx import Presentation
        return "\n".join(s.text_frame.text for s in Presentation(uploaded_file).slides
                         for shape in s.shapes if shape.has_text_frame)
    if ext in (".doc", ".docx"):
        from docx import Document
        return "\n".join(p.text for p in Document(uploaded_file).paragraphs)
    if ext in (".txt", ".md", ".markdown"):
        try:
            return uploaded_file.read().decode("utf-8")
        except UnicodeDecodeError:
            uploaded_file.seek(0)
            return uploaded_file.read().decode("gbk")
    # 图片类型提示
    if ext in (".jpg", ".jpeg", ".png", ".bmp", ".gif"):
        return "[IMAGE:" + uploaded_file.name + "]"
    return ""

def slice_text(text):
    lines, sections, cur_title, cur_lines = text.split("\n"), [], "引言", []
    hd = re.compile(r"^(#{1,4}\s+.+|\d+[\.\、\)]\s*.+|第[一二三四五六七八九十\d]+[章节].+)")
    for line in lines:
        s = line.strip()
        if hd.match(s) and len(s) < 80:
            if cur_lines:
                sections.append((cur_title, "\n".join(cur_lines).strip()))
            cur_title, cur_lines = s.lstrip("#").strip(), []
        else:
            cur_lines.append(line)
    if cur_lines:
        sections.append((cur_title, "\n".join(cur_lines).strip()))
    if len(sections) <= 1 and len(text) > 500:
        sections = [(f"片段{i+1}", p.strip()) for i, p in enumerate(text.split("\n\n")) if p.strip()]
    return sections

# ==================== 学科识别 ====================
DEFAULT_SUBJECTS = ["数据结构", "操作系统", "计算机网络", "计算机组成原理"]
KW_DICT = {
    "数据结构": ["树","图","栈","队列","链表","排序","查找","哈希","堆","遍历","二叉树"],
    "操作系统": ["进程","线程","死锁","内存","CPU","调度","文件系统","中断","信号量","内核"],
    "计算机网络": ["TCP","IP","HTTP","DNS","路由","协议","网络层","传输层","OSI","子网"],
    "计算机组成原理": ["CPU","指令","寄存器","ALU","总线","Cache","存储器","流水线","冯诺依曼"],
}

def get_subjects():
    subs = set(DEFAULT_SUBJECTS)
    for parent in [RAW, CPD]:
        if parent.exists():
            for d in parent.iterdir():
                if d.is_dir():
                    subs.add(d.name)
    return sorted(subs)

def classify_subject(text):
    existing = "、".join(get_subjects())
    r = call_ai(f"判断以下文本属于哪个学科。已知学科：{existing}；如果都不匹配，返回一个新学科名。只返回JSON：{{\"subject\":\"学科名\"}}\n\n{text[:3000]}",
                json_mode=True, temperature=0.3)
    try:
        s = json.loads(r).get("subject", "").strip()
        if s:
            for parent in [RAW, CPD]:
                (parent / s).mkdir(parents=True, exist_ok=True)
            return s
    except:
        pass
    scores = {s: sum(1 for k in ks if k in text) for s, ks in KW_DICT.items()}
    best = max(scores, key=scores.get) if any(scores.values()) else None
    return best if best and scores[best] > 0 else "其他学科"

def classify_entity(text):
    r = call_ai(f"判断以下笔记的类型（概念/算法/对比/考试技巧）并给出20字以内标题，返回JSON：{{\"type\":\"类型\",\"title\":\"标题\"}}\n\n{text[:2000]}",
                json_mode=True, temperature=0.3)
    try:
        d = json.loads(r)
        return d.get("type", "概念"), d.get("title", "学习笔记")
    except:
        if any(kw in text for kw in ["区别","对比","vs","不同于"]): return "对比", "概念对比"
        if any(kw in text for kw in ["步骤","流程","伪代码","算法"]): return "算法", "算法笔记"
        if any(kw in text for kw in ["技巧","记住","口诀","考点","易错"]): return "考试技巧", "考试技巧"
        return "概念", "概念理解"

# ==================== 知识收割 ====================
def save_raw(subject, title, content, src_type="文件上传", src_name=""):
    folder = RAW / subject
    folder.mkdir(parents=True, exist_ok=True)
    safe = re.sub(r"[\\/:*?\"<>|]", "_", title)[:50]
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    fp = folder / f"{ts}_{safe}.md"
    fm = f"---\ntype: {src_type}\nsubject: {subject}\nsource: {src_name or title}\ncreated: {datetime.now().strftime('%Y-%m-%d')}\n---\n"
    fp.write_text(fm + "\n" + content, encoding="utf-8")
    return fp

def extract_knowledge_points(course_id, text):
    """v5.0 新增：上传后 AI 自动提取结构化知识点"""
    prompt = f"""从以下文本中提取所有结构化知识点。返回严格JSON：
{{"knowledge_points":[{{"name":"知识点名称","description":"一句话描述","importance":"high/medium/low"}}]}}

文本：{text[:4000]}"""
    r = call_ai(prompt, json_mode=True, temperature=0.3)
    try:
        data = json.loads(r)
        count = 0
        for kp in data.get("knowledge_points", []):
            kp_upsert(course_id, kp.get("name",""), kp.get("description",""), kp.get("importance","medium"))
            count += 1
        return count
    except:
        return 0

def harvest_file(uploaded_file, course_id=None):
    text = parse_file(uploaded_file)
    if not text:
        return f"❌ 无法解析文件：{uploaded_file.name}"
    if text.startswith("[IMAGE:"):
        return f"⚠️ 检测到图片文件（{text[7:-1]}）。当前版本仅支持电子文档（PDF/PPTX/DOCX/TXT/MD）。图片识别功能计划在后续版本支持。"
    subject = classify_subject(text)
    sections = slice_text(text)
    saved = []
    for title, content in sections:
        if len(content.strip()) < 30:
            continue
        fp = save_raw(subject, title, content, src_type="文件上传", src_name=uploaded_file.name)
        saved.append(str(fp))
    # 提取知识点
    kp_count = 0
    if course_id:
        kp_count = extract_knowledge_points(course_id, text)
    kp_msg = f"，提取了 {kp_count} 个知识点" if kp_count else ""
    return f"✅ 已收割《{uploaded_file.name}》→【{subject}】{len(saved)}个片段{kp_msg}", saved, subject

def harvest_fragment(text, course_id=None):
    etype, title = classify_entity(text)
    subject = classify_subject(text)
    fp = save_raw(subject, title, text, src_type=etype, src_name="碎片输入")
    kp_count = 0
    if course_id:
        kp_count = extract_knowledge_points(course_id, text)
    kp_msg = f"，提取了 {kp_count} 个知识点" if kp_count else ""
    return f"✅ 碎片已收割：{title} →【{subject}】→【{etype}】{kp_msg}"

def list_raw(subject=None):
    dirs = [RAW / subject] if subject else [RAW / s for s in get_subjects()]
    files = []
    for d in dirs:
        if d.exists():
            files.extend(d.glob("*.md"))
    return files

def list_compiled(subject=None):
    dirs = [CPD / subject] if subject else [CPD / s for s in get_subjects()]
    files = []
    for d in dirs:
        if d.exists():
            files.extend(d.glob("*.md"))
    return files

def parse_fm(content):
    m = re.match(r"^---\s*\n(.*?)\n---\s*\n(.*)", content, re.DOTALL)
    if m:
        fm = {}
        for line in m.group(1).split("\n"):
            if ":" in line:
                k, _, v = line.partition(":")
                fm[k.strip()] = v.strip()
        return fm, m.group(2)
    return {}, content

# ==================== 知识编译 ====================
def compile_knowledge():
    raws = list_raw()
    if not raws:
        return "⚠️ 没有待编译的原始资料"
    cp_list = list_compiled()
    existing_str = "\n".join(f"- {p.stem}({p.parent.name})" for p in cp_list) if cp_list else "（无已有页面）"
    results, backlinks = [], {}

    for rf in raws:
        subject = rf.parent.name
        _, body = parse_fm(rf.read_text(encoding="utf-8"))
        prompt = f"""将以下学习资料编译为结构化知识页面。返回JSON：
{{"title":"页面标题(20字内)","subject":"学科","page_type":"概念/算法/对比/考试技巧","is_new":true,"merge_with":"","content":"Markdown内容(含[[页面名]]引用)","linked_pages":["被引用页面名"]}}

已有页面：{existing_str}

原始资料：{body[:5000]}"""
        r = call_ai(prompt, json_mode=True)
        try:
            d = json.loads(r)
        except:
            continue
        title = d.get("title", rf.stem)
        ps = d.get("subject", subject)
        safe = re.sub(r"[\\/:*?\"<>|]", "_", title)[:80]
        now = datetime.now().strftime("%Y-%m-%d")
        content = f"---\ntype: {d.get('page_type','概念')}\nsubject: {ps}\nsource: {rf.name}\ncreated: {now}\nupdated: {now}\n---\n\n{d.get('content', body)}"
        (CPD / ps).mkdir(parents=True, exist_ok=True)
        (CPD / ps / f"{safe}.md").write_text(content, encoding="utf-8")
        for lp in d.get("linked_pages", []):
            backlinks.setdefault(lp, []).append(safe)
        results.append(f"✅ {safe}（{ps}）")

    # 双向链接更新
    for page, refs in backlinks.items():
        for s in get_subjects():
            pp = CPD / s / f"{page}.md"
            if pp.exists():
                txt = pp.read_text(encoding="utf-8")
                txt = re.sub(r"\n## 被以下页面引用\n.*", "", txt, flags=re.DOTALL)
                txt += "\n## 被以下页面引用\n" + "\n".join(f"- [[{r}]]" for r in set(refs)) + "\n"
                pp.write_text(txt, encoding="utf-8")
                break

    # 全局索引
    idx_lines = [f"# 知识库全局索引\n\n> 更新于 {datetime.now().strftime('%Y-%m-%d %H:%M')}\n"]
    for s in get_subjects():
        idx_lines.append(f"\n## {s}\n")
        d = CPD / s
        if d.exists():
            for f in sorted(d.glob("*.md")):
                fm, _ = parse_fm(f.read_text(encoding="utf-8"))
                idx_lines.append(f"- [[{f.stem}]]  ({fm.get('type','')})")
        else:
            idx_lines.append("- （暂无页面）")
    IDX.write_text("\n".join(idx_lines) + "\n", encoding="utf-8")

    # 变更日志
    today = datetime.now().strftime("%Y-%m-%d")
    log_content = LOG.read_text(encoding="utf-8") if LOG.exists() else "# 知识库变更日志\n"
    if today not in log_content:
        log_content += f"\n## {today}\n"
    for entry in results:
        if entry not in log_content:
            log_content += f"- {entry}\n"
    LOG.write_text(log_content, encoding="utf-8")
    return "\n".join(results) if results else "✅ 所有资料已是最新"

# ==================== 知识问答（v5.0 多轮对话） ====================
def search_knowledge(query, top_k=5):
    qwords = set(query)
    scored = []
    for s in get_subjects():
        d = CPD / s
        if not d.exists():
            continue
        for f in d.glob("*.md"):
            content = f.read_text(encoding="utf-8")
            score = sum(content.count(w) for w in qwords)
            if any(w in f.stem for w in qwords):
                score += 10
            if score > 0:
                scored.append((f.stem, s, content, score))
    scored.sort(key=lambda x: -x[3])
    return scored[:top_k]

def answer_question(question, course_id=None):
    """v5.0：带多轮对话记忆 + SSE 流式"""
    hits = search_knowledge(question)
    if not hits:
        return "⚠️ 知识库中没有找到相关内容，请先上传资料并编译知识", []

    ctx = "\n\n---\n\n".join(f"【来源：{n}({s})】\n{c[:1500]}" for n, s, c, _ in hits)
    sources = [f"[[{n}]]（{s}）" for n, s, _, _ in hits]

    system_prompt = f"你是学习助手 SuperTutor。请基于以下资料回答学生问题。如资料不足请诚实告知，不要编造。\n\n参考资料：\n{ctx[:6000]}"

    # 构建 messages（含历史对话）
    messages_for_log = []
    if course_id:
        history = chat_recent(course_id, 10)
        history_text = "\n".join(f"[{r}]: {c[:200]}" for r, c in history)
        system_prompt += f"\n\n对话历史：\n{history_text}"
        messages_for_log = [(r, c) for r, c in history]

    # SSE 流式显示
    placeholder = st.empty()
    full_text = ""
    buffer = []
    for token in call_ai(question, system_prompt, stream=True):
        if isinstance(token, str) and token.startswith("⚠️"):
            placeholder.error(token)
            return token, sources
        buffer.append(token)
        full_text += token
        if len(buffer) >= 3 or token in "，。！？\n":
            placeholder.markdown(full_text + "▌")
            buffer = []
    placeholder.markdown(full_text)

    # 保存对话历史
    if course_id:
        chat_save(course_id, "user", question)
        chat_save(course_id, "assistant", full_text)

    return full_text, sources

# ==================== 薄弱点判断（4 条规则） ====================
def judge_weak(kp_id, accuracy, total, last_correct):
    if total < 3 and not last_correct:
        return False, "待观察"
    recent = db_query("""SELECT user_answer,correct_answer FROM quiz_question
        WHERE knowledge_point_id=? ORDER BY created_at DESC LIMIT 5""", (kp_id,), fetch=True)
    if len(recent) >= 2:
        cons_err = 0
        for ua, ca in recent:
            if ua and ua[0] == ca[0]:
                break
            cons_err += 1
        if cons_err >= 2:
            return True, f"连续错误{cons_err}次"
    if total >= 3 and accuracy < 0.6:
        return True, f"答题{total}次正确率{accuracy:.0%}<60%"
    if total >= 5 and accuracy >= 0.6 and len(recent) >= 3:
        last3 = sum(1 for ua, ca in recent[:3] if ua and ua[0] == ca[0]) / 3
        if last3 < 0.4:
            return True, f"总正确率{accuracy:.0%}但近3次降至{last3:.0%}"
    return False, ""

# ==================== 出题（v5.0：外键绑定知识点ID） ====================
def generate_quiz(subject, difficulty, count, course_id=None, weak_first=True):
    folder = CPD / subject
    knowledge_parts, page_names = [], []
    if folder.exists():
        for f in folder.glob("*.md"):
            content = f.read_text(encoding="utf-8")
            knowledge_parts.append(f"## {f.stem}\n{content[:2000]}")
            page_names.append(f.stem)
    if not knowledge_parts:
        return f"⚠️ {subject}暂无编译后的知识，请先进行知识编译", None

    # 构建知识点列表（v5.0：用 ID 不用名称字符串）
    kp_rows = kp_list(course_id) if course_id else []
    kp_map = {row[2]: row[0] for row in kp_rows}  # name → id
    kp_list_str = "\n".join(
        f"ID={row[0]}, 名称={row[2]}, 掌握度={row[5]:.0%}" + ("（薄弱⚠️）" if row[8] else "")
        for row in kp_rows
    ) if kp_rows else "（暂无知识点数据）"

    weak_kps = [row[2] for row in kp_rows if row[8]] if weak_first else []
    diff_map = {"自动匹配": "基础30%+中等50%+困难20%", "基础": "全部基础", "中等": "全部中等", "困难": "全部困难"}
    diff_inst = diff_map.get(difficulty, "基础30%+中等50%+困难20%")
    weak_inst = f"优先针对以下薄弱知识点出题：{', '.join(weak_kps[:5])}" if weak_kps else ""

    prompt = f"""基于以下知识点列表和知识内容生成{count}道选择题(4选项A/B/C/D)。{diff_inst}。{weak_inst}
每道题必须包含 knowledge_point_id 字段（用知识点ID，不是名称）。返回严格JSON：
{{"questions":[{{"knowledge_point_id":1,"question":"题目","options":["A.xx","B.xx","C.xx","D.xx"],"answer":"A","explanation":"解析(100字内)","source_page":"页面名","difficulty":"基础/中等/困难"}}]}}

知识点列表（出题时用 knowledge_point_id 关联）：
{kp_list_str}

可用页面：{', '.join(page_names)}

知识内容：\n{''.join(knowledge_parts)[:8000]}"""

    r = call_ai(prompt, json_mode=True, temperature=0.9)
    try:
        questions = json.loads(r).get("questions", [])
        # 补齐 knowledge_point_id：如果 AI 没返回或返回的名称不在映射中，尝试按 source_page 匹配
        for q in questions:
            if not q.get("knowledge_point_id"):
                src = q.get("source_page", "")
                for name, kid in kp_map.items():
                    if name in src or src in name:
                        q["knowledge_point_id"] = kid
                        break
        return questions, page_names
    except:
        return f"⚠️ 题目解析失败：{r[:200]}", None

# ==================== 复习计划（v5.0：艾宾浩斯遗忘曲线） ====================
def generate_plan(course_id, exam_date, daily_minutes):
    days = (exam_date - datetime.now().date()).days
    if days <= 0:
        return "⚠️ 考试日期已过，请重新设置"

    # 取课程掌握度数据
    kp_rows = kp_list(course_id)
    weak_rows = kp_get_weak(course_id)
    weak_str = "\n".join(
        f"- {row[2]}（正确率{row[5]:.0%}，{row[9]}）"
        for row in weak_rows
    ) if weak_rows else "（暂无薄弱点）"

    # 学科统计
    stats = []
    for s in get_subjects():
        d = CPD / s
        pc = len(list(d.glob("*.md"))) if d.exists() else 0
        stats.append(f"- {s}：{pc}个知识页面")
    for row in kp_rows:
        found = False
        for i, st in enumerate(stats):
            if row[2] in st:
                found = True
                break
        if not found:
            stats.append(f"- {row[2]}：掌握度{row[5]:.0%}" + ("⚠️薄弱" if row[8] else ""))

    prompt = f"""制定期末复习计划。遵循以下规则：

1. 艾宾浩斯遗忘曲线：每个知识点首次学习后，在第 1、2、4、7、15 天后需安排复习
2. 薄弱点（正确率<60%）优先安排在前期
3. 每天不超过 {daily_minutes} 分钟
4. 每天不超过 2 个学科，避免频繁切换
5. 最后 3 天安排总复习 + 所有错题回顾

考试日期：{exam_date}（还有 {days} 天）
每日可用时长：{daily_minutes} 分钟

学科统计：
{chr(10).join(stats)}

薄弱点（优先安排）：
{weak_str}

按日期列出每日任务（学科 + 知识点 + 预计用时 + 是否做题）。用 Markdown 表格格式。"""

    return call_ai(prompt)

# ==================== Streamlit UI ====================
def main():
    st.set_page_config(page_title="SuperTutor v5.0", page_icon="🎓", layout="wide")

    # Session state 初始化
    for k, v in {"api_key": "", "current_course_id": None, "quiz_qs": [], "quiz_done": False,
                  "quiz_res": [], "quiz_record_id": None, "plan": "", "review_mode": None}.items():
        if k not in st.session_state:
            st.session_state[k] = v

    init_db()

    # 确保默认课程存在
    courses = course_list()
    if not courses:
        cid = course_create("默认课程", "首次使用自动创建的课程")
        st.session_state.current_course_id = cid
    elif st.session_state.current_course_id is None:
        st.session_state.current_course_id = courses[0][0]

    # 确保知识库目录存在
    for s in get_subjects():
        (RAW / s).mkdir(parents=True, exist_ok=True)
        (CPD / s).mkdir(parents=True, exist_ok=True)
    if not IDX.exists():
        IDX.write_text("# 知识库全局索引\n\n（暂无内容）\n", encoding="utf-8")
    if not LOG.exists():
        LOG.write_text("# 知识库变更日志\n", encoding="utf-8")

    cid = st.session_state.current_course_id
    course = course_get(cid)

    # ===== 侧边栏 =====
    with st.sidebar:
        st.title("🎓 SuperTutor v5.0")

        # API Key
        key = st.text_input("🔑 API Key", value=st.session_state.api_key, type="password",
                            placeholder="sk-... 或在 config.properties 中配置")
        if key != st.session_state.api_key:
            st.session_state.api_key = key

        st.divider()

        # 课程管理
        st.subheader("📚 课程")
        for c in courses:
            c_name = c[1]
            c_exam = f" 📅{c[3]}" if c[3] else ""
            if st.button(f"{'📍' if c[0] == cid else '  '} {c_name}{c_exam}",
                         key=f"course_{c[0]}", use_container_width=True,
                         help=f"切换到{c_name}"):
                st.session_state.current_course_id = c[0]
                st.session_state.quiz_qs = []
                st.session_state.quiz_done = False
                st.rerun()

        # 新建课程
        with st.expander("➕ 新建课程"):
            new_name = st.text_input("课程名称", key="new_course_name")
            new_desc = st.text_input("描述（可选）", key="new_course_desc")
            if st.button("创建课程") and new_name:
                new_id = course_create(new_name, new_desc)
                st.session_state.current_course_id = new_id
                st.rerun()

        # 删除课程（非默认课程才可删）
        if courses and len(courses) > 1 and cid != courses[0][0]:
            if st.button("🗑️ 删除当前课程", type="secondary", use_container_width=True):
                course_delete(cid)
                st.session_state.current_course_id = course_list()[0][0]
                st.rerun()

        st.divider()

        # 当前课程知识点
        if course:
            kp_rows = kp_list(cid)
            weak_count = sum(1 for r in kp_rows if r[7])
            st.subheader(f"📊 {course[1]}")
            st.metric("知识点", f"{len(kp_rows)}个")
            st.metric("薄弱", f"{weak_count}个")
            if kp_rows:
                with st.expander("知识点详情"):
                    for row in kp_rows:
                        kp_id, _, name, desc, imp, mastery, tot, corr, iw, reason, lr = row
                        color = "🟢" if mastery >= 0.8 else ("🟡" if mastery >= 0.5 else "🔴")
                        st.markdown(f"{color} **{name}** {mastery:.0%}" + (" ⚠️" if iw else ""))
                        if reason:
                            st.caption(f"  {reason}")

        st.divider()

        # 操作按钮
        if st.button("🔄 知识编译", use_container_width=True):
            with st.spinner("编译中..."):
                st.success(compile_knowledge())
                st.rerun()

        # 使用说明
        with st.expander("📖 使用说明"):
            st.markdown("**1. 新建课程** → **2. 上传资料** → **3. 知识编译** → **4. 刷题/问答/计划**")
            st.caption("上传资料后 AI 自动提取知识点到侧边栏")

    # ===== 主标签页 =====
    t1, t2, t3 = st.tabs(["📚 知识收割", "✏️ 精准刷题", "📅 复习规划"])

    # ---- 标签页1：知识收割 ----
    with t1:
        total_pages = sum(1 for _ in list_compiled())
        kp_rows_now = kp_list(cid)
        mastered = sum(1 for r in kp_rows_now if r[5] >= 0.7 and r[6] >= 3)
        weak_n = sum(1 for r in kp_rows_now if r[7])
        days_left = ""
        if course and course[3]:
            delta = (datetime.strptime(course[3], "%Y-%m-%d").date() - datetime.now().date()).days
            days_left = str(max(0, delta))

        c1, c2, c3, c4 = st.columns(4)
        c1.metric("📄 知识页面", total_pages)
        c2.metric("🧠 知识点", len(kp_rows_now))
        c3.metric("🔴 薄弱", weak_n)
        c4.metric("⏰ 距考试", f"{days_left}天" if days_left else "未设置")

        st.divider()
        st.subheader("📤 上传资料（PDF/PPT/Word/TXT/MD）")
        files = st.file_uploader("选择文件", ["pdf","ppt","pptx","doc","docx","txt","md"],
                                 accept_multiple_files=True, label_visibility="collapsed")
        if files:
            for f in files:
                with st.spinner(f"处理 {f.name}..."):
                    result = harvest_file(f, cid)
                    msg = result[0] if isinstance(result, tuple) else result
                    if "✅" in msg:
                        st.success(msg)
                    else:
                        st.warning(msg)

        st.divider()
        st.subheader("✍️ 碎片知识收割")
        frag = st.text_area("输入易错点或顿悟笔记",
                            placeholder="例如：二叉树中序遍历是左-根-右，画图时从根向左走到尽头再回溯...",
                            height=100, label_visibility="collapsed")
        if st.button("💾 收割碎片", key="btn_frag"):
            if frag.strip():
                with st.spinner("识别中..."):
                    st.success(harvest_fragment(frag, cid))
            else:
                st.warning("请输入内容")

        st.divider()
        st.subheader("❓ 知识问答（多轮对话）")
        q = st.text_input("输入问题", placeholder="例如：进程和线程的区别是什么？",
                          label_visibility="collapsed", key="qa_input")
        if st.button("🔍 提问", key="btn_qa"):
            if q.strip():
                if not list_compiled():
                    st.warning("知识库为空，请先上传资料并编译")
                else:
                    with st.spinner("思考中..."):
                        ans, srcs = answer_question(q, cid)
                    if srcs:
                        st.caption("📎 参考来源：" + " | ".join(srcs))

    # ---- 标签页2：精准刷题 ----
    with t2:
        if sum(1 for _ in list_compiled()) == 0:
            st.info("知识库为空，请先在「知识收割」页上传资料并在侧边栏编译")
        else:
            c1, c2, c3, c4 = st.columns(4)
            with c1:
                q_subject = st.selectbox("学科", get_subjects(), key="qs")
            with c2:
                q_diff = st.selectbox("难度", DIFFS, key="qd")
            with c3:
                q_count = st.selectbox("题量", COUNTS, key="qc")
            with c4:
                kp_data = kp_list(cid)
                has_weak = any(r[7] for r in kp_data) if kp_data else False
                weak_first = st.checkbox("薄弱点优先", value=True, disabled=not has_weak,
                                          help="有薄弱知识点时建议开启")

            if st.button("🎲 生成题目", type="primary", use_container_width=True):
                with st.spinner(f"生成{q_count}道题..."):
                    qs, _ = generate_quiz(q_subject, q_diff, q_count, cid, weak_first)
                    if isinstance(qs, str):
                        st.error(qs)
                        st.session_state.quiz_qs = []
                    else:
                        st.session_state.quiz_qs = qs or []
                        st.session_state.quiz_done = False
                        st.session_state.quiz_res = []
                        st.rerun()

            qs = st.session_state.quiz_qs
            if qs:
                st.subheader(f"📝 {len(qs)} 道题")
                answers = {}
                for i, qz in enumerate(qs):
                    st.markdown(f"**{i+1}. {qz['question']}**")
                    st.caption(f"{qz.get('difficulty','基础')} | 来源：{qz.get('source_page','未知')}")
                    choice = st.radio(f"q{i}", qz["options"], key=f"qr_{i}", index=None,
                                      disabled=st.session_state.quiz_done, label_visibility="collapsed")
                    if choice:
                        answers[i] = choice
                    st.divider()

                if not st.session_state.quiz_done:
                    if st.button("✅ 提交答案", type="primary", use_container_width=True):
                        if len(answers) < len(qs):
                            st.warning(f"还有{len(qs)-len(answers)}题未答")
                        else:
                            # 创建 quiz_record（INSERT + last_insert_rowid 必须在同一连接中）
                            qr_conn = sqlite3.connect(str(DB))
                            qr_conn.execute("PRAGMA foreign_keys = ON")
                            qr_conn.execute("INSERT INTO quiz_record(course_id,quiz_type,total_questions) VALUES(?,'periodic',?)",
                                           (cid, len(qs)))
                            qr_id = qr_conn.execute("SELECT last_insert_rowid()").fetchone()[0]
                            qr_conn.commit()
                            qr_conn.close()
                            st.session_state.quiz_record_id = qr_id

                            res = []
                            correct_n = 0
                            for i, qz in enumerate(qs):
                                ua = answers[i]
                                ul = ua[0] if ua else ""
                                cl = qz["answer"].strip()[0]
                                ok = ul == cl
                                if ok:
                                    correct_n += 1

                                kp_id = qz.get("knowledge_point_id")
                                # 验证 kp_id 是否真实存在于 knowledge_point 表
                                valid_kp_ids = {r[0] for r in kp_list(cid)}
                                if kp_id not in valid_kp_ids:
                                    kp_id = None  # AI 幻觉的 ID，置空
                                # 如果 kp_id 为空，尝试按 source_page 匹配
                                if not kp_id:
                                    src = qz.get("source_page", "")
                                    for row in kp_list(cid):
                                        if row[2] in src or src in row[2]:
                                            kp_id = row[0]
                                            break
                                # 保存每道题
                                db_query("""INSERT INTO quiz_question(quiz_record_id,course_id,knowledge_point_id,
                                    question_index,question_type,content,options,correct_answer,user_answer,explanation,difficulty)
                                    VALUES(?,?,?,?,?,?,?,?,?,?,?)""",
                                    (qr_id, cid, kp_id, i+1, "choice", qz["question"],
                                     json.dumps(qz.get("options",[]), ensure_ascii=False),
                                     qz["answer"], ua, qz.get("explanation",""), qz.get("difficulty","基础")))

                                # 更新知识点掌握度
                                if kp_id:
                                    kp_update_mastery(kp_id, ok)
                                else:
                                    # 回退：按 source_page 匹配知识点名称
                                    src = qz.get("source_page", "")
                                    for row in kp_list(cid):
                                        if row[2] in src or src in row[2]:
                                            kp_update_mastery(row[0], ok)
                                            break

                                # 答错 → 记录到 wrong_answer
                                if not ok:
                                    qq_id = db_query("SELECT last_insert_rowid()", fetch=True)[0][0]
                                    wrong_record(cid, kp_id, qq_id, "choice", qz["question"],
                                                 json.dumps(qz.get("options",[]), ensure_ascii=False),
                                                 qz["answer"], ua, qz.get("explanation",""))

                                res.append({"i": i, "q": qz["question"], "ua": ua, "ca": qz["answer"],
                                             "ok": ok, "exp": qz.get("explanation",""),
                                             "src": qz.get("source_page","未知"),
                                             "diff": qz.get("difficulty","基础")})

                            # 更新 quiz_record 的统计
                            acc = correct_n / len(qs) if qs else 0
                            db_query("UPDATE quiz_record SET correct_count=?, accuracy=? WHERE id=?",
                                     (correct_n, acc, qr_id))

                            st.session_state.quiz_res = res
                            st.session_state.quiz_done = True
                            st.rerun()
                else:
                    res = st.session_state.quiz_res
                    cc = sum(1 for r in res if r["ok"])
                    pct = cc / len(res) * 100 if res else 0
                    color = "green" if pct >= 80 else ("orange" if pct >= 60 else "red")
                    st.markdown(f"### 成绩：<span style='color:{color}'>{cc}/{len(res)}（{pct:.0f}%）</span>",
                                unsafe_allow_html=True)
                    for r in res:
                        icon = "✅" if r["ok"] else "❌"
                        st.markdown(f"{icon} **{r['i']+1}. {r['q']}**")
                        if r["ok"]:
                            st.success(f"你的答案：{r['ua']} ✔️")
                        else:
                            st.error(f"你的答案：{r['ua']} | 正确答案：**{r['ca']}**")
                        with st.expander("查看解析"):
                            st.markdown(r["exp"])
                            st.caption(f"来源：{r['src']}")
                    if st.button("🔄 再来一组", use_container_width=True):
                        st.session_state.quiz_qs = []
                        st.session_state.quiz_done = False
                        st.session_state.quiz_res = []
                        st.rerun()

            # ---- 错题本（v5.0：支持复习重做） ----
            st.divider()
            st.subheader("📕 错题本")
            wrong_rows = wrong_list_unmastered(cid)
            if not wrong_rows:
                st.success("🎉 暂无未掌握错题！")
            else:
                st.caption(f"共 {len(wrong_rows)} 道未掌握错题")
                for row in wrong_rows:
                    wid, w_cid, kp_id, qq_id, qtype, content, opts, correct_ans, user_ans, expl, rc, lr, mastered, created = row
                    # 找关联知识点名称
                    kp_name = ""
                    kp_rows_map = {r[0]: r[2] for r in kp_list(cid)}
                    if kp_id and kp_id in kp_rows_map:
                        kp_name = kp_rows_map[kp_id]

                    kp_label = f" [{kp_name}]" if kp_name else ""
                    reviewed_label = f" | 复习{rc}次" if rc else ""
                    with st.expander(f"❌ {content[:60]}...{kp_label} | 选了：{user_ans} → 正确答案：{correct_ans}{reviewed_label}"):
                        # 复习模式
                        if st.session_state.get("review_mode") == wid:
                            st.markdown(f"**题目**：{content}")
                            try:
                                opts_list = json.loads(opts) if isinstance(opts, str) else opts
                            except:
                                opts_list = []
                            re_answer = st.radio("重新选择", opts_list, key=f"review_{wid}", index=None)
                            if st.button("✅ 确认复习答案", key=f"confirm_{wid}"):
                                if re_answer:
                                    re_ok = re_answer[0] == correct_ans[0]
                                    mastered = wrong_review(wid, re_answer, re_ok)
                                    if re_ok:
                                        st.success("🎉 回答正确！已标记为已掌握")
                                        # 同时更新知识点掌握度
                                        if kp_id:
                                            kp_update_mastery(kp_id, True)
                                        st.session_state.review_mode = None
                                        st.rerun()
                                    else:
                                        st.error(f"回答错误。正确答案：{correct_ans}")
                                        st.markdown(f"**解析**：{expl}")
                                        if st.button("知道了", key=f"ok_{wid}"):
                                            st.session_state.review_mode = None
                                            st.rerun()
                            if st.button("取消", key=f"cancel_{wid}"):
                                st.session_state.review_mode = None
                                st.rerun()
                        else:
                            st.markdown(f"**题目**：{content}")
                            st.markdown(f"**你的答案**：{user_ans}")
                            st.markdown(f"**正确答案**：{correct_ans}")
                            st.markdown(f"**解析**：{expl}")
                            if st.button("🔄 重新作答", key=f"redo_{wid}"):
                                st.session_state.review_mode = wid
                                st.rerun()

    # ---- 标签页3：复习规划 ----
    with t3:
        st.subheader("📅 复习规划")
        if not course:
            st.warning("请先在侧边栏创建课程")
        else:
            c1, c2 = st.columns(2)
            with c1:
                default_date = datetime.strptime(course[3], "%Y-%m-%d").date() if course[3] else datetime.now().date() + timedelta(days=30)
                ed = st.date_input("考试日期", value=default_date, min_value=datetime.now().date(), key="ed_input")
                if str(ed) != course[3]:
                    course_update(cid, exam_date=ed)
            with c2:
                dh = st.number_input("每日学习时长（分钟）", 30, 960,
                                     value=course[4] if course[4] else 240, step=30, key="dh_input")
                if dh != course[4]:
                    course_update(cid, daily_minutes=dh)

            # 掌握度总览
            st.subheader("📊 掌握度总览")
            kp_data = kp_list(cid)
            if kp_data:
                for row in kp_data:
                    kp_id, _, name, desc, imp, mastery, tot, corr, iw, reason, lr = row
                    col1, col2 = st.columns([4, 1])
                    with col1:
                        color = "🟢" if mastery >= 0.8 else ("🟡" if mastery >= 0.5 else "🔴")
                        label = f"{color} {name}" + (" ⚠️" if iw else "")
                        st.progress(min(mastery, 1.0), text=f"{label} {mastery:.0%}" if tot > 0 else f"{color} {name} 未练习")
                    with col2:
                        st.caption(f"{tot}题" if tot else "新")
            else:
                st.info("暂无知识点数据，请先上传资料")

            st.divider()

            # 薄弱点
            st.subheader("🔴 薄弱点与待观察")
            weak_rows = kp_get_weak(cid)
            obs_rows = kp_get_observing(cid)
            if weak_rows:
                for row in weak_rows:
                    st.markdown(f"⚠️ **{row[2]}** — 正确率{row[5]:.0%}，{row[9]}")
            if obs_rows:
                for row in obs_rows:
                    st.markdown(f"👀 **{row[2]}** — {row[9]}")
            if not weak_rows and not obs_rows:
                st.success("暂无薄弱点，继续保持！")

            st.divider()

            # 学习历史趋势
            st.subheader("📈 测验历史")
            qr_rows = db_query("SELECT * FROM quiz_record WHERE course_id=? ORDER BY created_at DESC LIMIT 10",
                               (cid,), fetch=True)
            if qr_rows:
                for qr in qr_rows:
                    qr_id, _, qtype, total, correct, acc, created = qr
                    st.markdown(f"{created[:10]} | {qtype} | {correct}/{total}（{acc:.0%}）")
                # 趋势
                accs = [qr[5] for qr in qr_rows if qr[5] is not None]
                if len(accs) >= 2:
                    trend = "上升 📈" if accs[1] > accs[0] else ("下降 📉" if accs[1] < accs[0] else "持平 ➡️")
                    st.caption(f"最近两次测验趋势：{trend}（{accs[1]:.0%} → {accs[0]:.0%}）")
            else:
                st.info("暂无测验记录")

            st.divider()

            # 复习计划
            if st.button("📋 生成复习计划", type="primary", use_container_width=True):
                if not list_compiled():
                    st.warning("知识库为空，请先上传资料并编译")
                elif not ed:
                    st.warning("请先设置考试日期")
                else:
                    with st.spinner("生成个性化复习计划（含艾宾浩斯遗忘曲线）..."):
                        st.session_state.plan = generate_plan(cid, ed, dh)
            if st.session_state.plan:
                st.markdown("### 📋 每日复习计划")
                st.markdown(st.session_state.plan)
                st.download_button("📥 下载复习计划", st.session_state.plan,
                                   f"复习计划_{ed}.md", "text/markdown")

if __name__ == "__main__":
    main()
