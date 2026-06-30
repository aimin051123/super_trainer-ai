"""文件解析：PDF/PPTX/DOCX/TXT/MD + 图片 OCR"""
from pathlib import Path

_ocr_reader = None

def _get_ocr():
    """懒加载 easyocr Reader，首次调用时初始化"""
    global _ocr_reader
    if _ocr_reader is None:
        import easyocr
        _ocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=False, verbose=False)
    return _ocr_reader

def parse_image(uploaded_file):
    """使用 easyocr 识别图片中的文字"""
    try:
        import tempfile, os
        reader = _get_ocr()
        suffix = Path(uploaded_file.name).suffix or ".png"
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(uploaded_file.read())
            tmp_path = tmp.name
        results = reader.readtext(tmp_path, detail=0)
        os.unlink(tmp_path)
        return "\n".join(results) if results else ""
    except Exception as e:
        return f"[OCR 失败: {e}]"

def parse_file(uploaded_file):
    name = uploaded_file.name.lower()
    ext = Path(name).suffix
    if ext == ".pdf":
        from pypdf import PdfReader
        text = "\n".join(p.extract_text() or "" for p in PdfReader(uploaded_file).pages)
        if text.strip():
            return text
        # PDF 提取为空，可能是扫描件，尝试 OCR
        uploaded_file.seek(0)
        return parse_image(uploaded_file)
    if ext in (".ppt", ".pptx"):
        from pptx import Presentation
        return "\n".join(s.text_frame.text for s in Presentation(uploaded_file).slides
                         for shape in s.shapes if shape.has_text_frame)
    if ext in (".doc", ".docx"):
        from docx import Document
        return "\n".join(p.text for p in Document(uploaded_file).paragraphs)
    if ext in (".txt", ".md", ".markdown"):
        try:
            return uploaded_file.read().decode("utf-8")
        except UnicodeDecodeError:
            uploaded_file.seek(0)
            return uploaded_file.read().decode("gbk")
    if ext in (".jpg", ".jpeg", ".png", ".bmp", ".gif"):
        return parse_image(uploaded_file)
    return ""
